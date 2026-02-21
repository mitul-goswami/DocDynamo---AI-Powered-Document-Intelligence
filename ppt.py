import re
import json
import base64
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
# Initialize Mistral client with API key
from pathlib import Path
from mistralai import DocumentURLChunk, ImageURLChunk, TextChunk
import json
from groq import Groq
from mistralai import Mistral
from dotenv import load_dotenv
import os
import ast
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
# print(GROQ_API_KEY)

def convert_dict(txt):
    """
    Parse different input formats to extract title and content.
    Handles various JSON and string input formats.
    """
    # Remove outer quotes and strip whitespace
    txt = txt.strip('"').strip()
    org_txt = txt
    txt = txt[txt.find("[") + 1:txt.find("]") + 1]
    if len(txt) < 20:
        txt = org_txt
    try:
        # Try parsing as a JSON-like array
        if txt.startswith("["):
            # Split the array into title and content
            parts = txt.strip('[]').split(',', 1)
            title = parts[0].strip().strip('"')
            content = parts[1].strip().strip('"""')
            return {
                "title": title,
                "content": content
            }
        
        # If direct dictionary-like input
        if isinstance(txt, str):
            # Try to extract title and content
            title_match = re.search(r'"(.*?)"', txt)
            content_match = re.search(r'"""(.*?)"""', txt, re.DOTALL)
            
            title = title_match.group(1) if title_match else "Untitled"
            content = content_match.group(1).strip() if content_match else txt
            
            return {
                "title": title,
                "content": content
            }
    except Exception as e:
        print(f"Error parsing input: {e}")
        return {
            "title": "Error",
            "content": txt
        }

def parse_markdown_line(line):
    """
    Parse a markdown bullet line formatted as:
    "- **Heading**: Content"
    Returns a tuple (heading, content). If no heading is found, returns (None, line).
    """
    line = line.strip()
    if line.startswith("-"):
        line = line[1:].strip()
    pattern = r"\*\*(.*?)\*\*:\s*(.*)"
    match = re.match(pattern, line)
    if match:
        return match.group(1), match.group(2)
    else:
        return None, line

# def create_presentation(data_array, file_path):
#     """Create PPT with adaptive layout for images"""
#     prs = Presentation()
    
#     for slide_entry in data_array:
#         content_str, images = slide_entry[0], slide_entry[1]
#         slide_data = convert_dict(content_str)
#         # print(slide_data)
#         has_images = len(images) > 0

#         # Choose layout based on image presence
#         if has_images:
#             slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
#         else:
#             slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title+Content layout

#         # Set title (different positioning for image slides)
#         if has_images:
#             # Create a custom textbox for the title
#             title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
#             title_tf = title_box.text_frame
#         else:
#             title_tf = slide.shapes.title.text_frame

#         # Set slide title
#         title_tf.text = slide_data['title']
#         title_paragraph = title_tf.paragraphs[0]
#         title_paragraph.font.size = Pt(32)
#         title_paragraph.font.bold = True
#         title_paragraph.font.name = 'Arial'
#         title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
#         title_paragraph.alignment = 1  # Center-align

#         # Set content parameters based on image presence
#         # content_font_size = Pt(14 if has_images else 16)
#         content_width = Inches(5.5) if has_images else Inches(12)
#         content_left = Inches(0.5) if has_images else Inches(1)

#         # Create content textbox or use placeholder
#         if has_images:
#             content_box = slide.shapes.add_textbox(content_left, Inches(1.5), content_width, Inches(5))
#             content_tf = content_box.text_frame
#         else:
#             content_tf = slide.placeholders[1].text_frame
#             content_tf.clear()

#         # Process each line of content
#         for line in slide_data['content'].split('\n'):
#             if line.strip() == "":
#                 continue
            
#             heading, text = parse_markdown_line(line)
#             p = content_tf.add_paragraph()
            
#             # Add bold heading if exists
#             if heading:
#                 run_heading = p.add_run()
#                 run_heading.text = f"{heading}: "
#                 run_heading.font.bold = True
            
#             # Add content text
#             run_text = p.add_run()
#             run_text.text = text
            
#             # Style paragraph runs
#             for run in p.runs:
#                 run.font.size = Pt(14)
#                 run.font.name = 'Calibri'
#                 run.font.color.rgb = RGBColor(51, 51, 51)
#             p.alignment = 1  # Center-align
#             p.space_after = Pt(5)  # Use Pt for spacing

#         # Add images if present
#         if has_images:
#             img_left = Inches(3)
#             img_top = Inches(4)
#             max_width = Inches(4)
            
#             for img_data in images:
#                 try:
#                     # Handle base64 data URI
#                     if 'base64,' in img_data:
#                         img_data = img_data.split(',', 1)[1]
                    
#                     image_bytes = base64.b64decode(img_data)
#                     image_stream = BytesIO(image_bytes)
                    
#                     # Add image to slide
#                     pic = slide.shapes.add_picture(image_stream, img_left, img_top, width=max_width)
#                     # Update position for next image
#                     img_top += pic.height + Inches(0.1)
#                 except Exception as e:
#                     print(f"Error processing image: {e}")

#     # Save presentation
#     prs.save(file_path)
#     print(f"✅ Professional presentation saved to {file_path}")
def create_presentation(data_array, file_path):
    """Create PPT with adaptive layout for images and prevent text overflow."""
    prs = Presentation()
    
    for slide_entry in data_array:
        content_str, images = slide_entry[0], slide_entry[1]
        slide_data = convert_dict(content_str)
        print(slide_data)
        has_images = len(images) > 0

        # Choose layout based on image presence
        if has_images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title+Content layout

        # -- TITLE SETUP --
        if has_images:
            # Create a custom textbox for the title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8), Inches(0.8))
            title_tf = title_box.text_frame
        else:
            title_tf = slide.shapes.title.text_frame

        # Enable word wrap and auto-size to prevent overflow
        title_tf.word_wrap = True 
        # Set slide title
        title_tf.text = slide_data['title']
        title_paragraph = title_tf.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.name = 'Arial'
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102) # Center-align

        # -- CONTENT SETUP --
        content_width = Inches(12)
        content_left = Inches(0.5) if has_images else Inches(1)
        
        # Create content textbox or use placeholder
        if has_images:
            content_box = slide.shapes.add_textbox(content_left, Inches(1.5), content_width, Inches(5))
            content_tf = content_box.text_frame
        else:
            content_tf = slide.placeholders[1].text_frame
            content_tf.clear()

        # Enable word wrap and auto-size for content
        content_tf.word_wrap = True
        
        # Optional: add small margins so text doesn't touch edges
        content_tf.margin_left = Inches(0.1)
        content_tf.margin_right = Inches(0.1)
        content_tf.margin_top = Inches(0.1)
        content_tf.margin_bottom = Inches(0.1)

        # Process each line of content
        for line in slide_data['content'].split('\n'):
            if line.strip() == "":
                continue
            
            heading, text = parse_markdown_line(line)
            p = content_tf.add_paragraph()
            
            # Add bold heading if exists
            if heading:
                run_heading = p.add_run()
                run_heading.text = f"{heading}: "
                run_heading.font.bold = True
            
            # Add content text
            run_text = p.add_run()
            run_text.text = text
            
            # Style paragraph runs
            for run in p.runs:
                run.font.size = Pt(14)
                run.font.name = 'Calibri'
                run.font.color.rgb = RGBColor(51, 51, 51)
                
            p.space_after = Pt(5)

        # -- IMAGES (unchanged) --
        if has_images:
            img_left = Inches(4)
            img_top = Inches(5)
            max_width = Inches(3)
            
            for img_data in images:
                try:
                    # Handle base64 data URI
                    if 'base64,' in img_data:
                        img_data = img_data.split(',', 1)[1]
                    
                    image_bytes = base64.b64decode(img_data)
                    image_stream = BytesIO(image_bytes)
                    
                    # Add image to slide
                    pic = slide.shapes.add_picture(image_stream, img_left, img_top, width=max_width)
                    # Update position for next image
                    img_top += pic.height + Inches(0.1)
                except Exception as e:
                    print(f"Error processing image: {e}")

    # Save presentation
    prs.save(file_path)
    print(f"✅ Professional presentation saved to {file_path}")
# pdf_bytes = ""
# Verify PDF file exists
def get_pdf_path(path_bytes):
    with open('test.pdf',"wb") as f:
        f.write(path_bytes)
        print("success")
def pdf_to_ppt(client_g):
    api_key = "ulAo0oH7uXYxLtDh56yQuQKjTiZwzws5" # Replace with your API key
    client = Mistral(api_key=api_key)
    
    def image_to_detail(base64_data_url):
        chat_response = client.chat.complete(
            model="pixtral-12b-latest",
            messages=[
                {
                    "role": "user",
                    "content": [
                        ImageURLChunk(image_url=base64_data_url),
                        TextChunk(
                            text=(
                                """Please provide a detailed yet concise description of the image in JSON format. Your JSON response should include the following keys:

        what: Identify the main subjects or elements present in the image.
        why: Explain the possible intent, context, or significance behind these elements.

        Ensure your response strictly follows JSON format without additional commentary."""
                                "The output should be strictly be json with no extra commentary"
                            )
                        ),
                    ],
                }
            ],
            response_format={"type": "json_object"},
            temperature=0,
        )

        response_dict = json.loads(chat_response.choices[0].message.content)
        return response_dict
    
    def summarize_data_for_ppt(data):
          chat_completion = client_g.chat.completions.create(
              messages=[
                  {
                      "role": "system",
                      "content": """
  You're a helpful assistant that transforms input data into an **array** with exactly **two elements**: "title" and "content". Internally, use chain-of-thought reasoning to extract 8 to 10 key points without revealing your process. Retain the provided title and produce a concise summary in **Markdown format**, strictly under **25 words**. For any flowcharts or diagrams, convert their meaning into narrative text without referencing visuals.

  **Response Format Example:**

json
  ["Understanding Generative AI Techniques",
    '''- **Style Transfer**: Applies art style to images.
  - **Data Augmentation**: Creates synthetic data for improved model training.
  - **Transfer Learning**: Uses pre-trained models for specific tasks.
  - **Deepfakes**: Synthetic media created using AI to deceive.
  - **Autoencoder**: Compresses data for reconstruction.
  - **GAN Loss Functions**: Guides stable generator and discriminator training.
  - **Attention Mechanism**: Enhances contextual dependency in NLP text generation.''']
"""
                  },
                  {
                      "role": "user",
                      "content": f"{data}",
                  }
              ],
              model="llama-3.1-8b-instant",
          )
          return chat_completion.choices[0].message.content
    
    # Verify PDF file exists
    pdf_file = Path("test.pdf")
    # assert file_path.is_file()

    # Upload PDF file 
    uploaded_file = client.files.upload(
        file={
            "file_name": pdf_file.stem,
            "content": pdf_file.read_bytes(),
        },
        purpose="ocr",
    )
    print("success_2")
    # Get URL for the uploaded file
    signed_url = client.files.get_signed_url(file_id=uploaded_file.id, expiry=1)

    # Process PDF with OCR, including embedded images
    pdf_response = client.ocr.process(
        document=DocumentURLChunk(document_url=signed_url.url),
        model="mistral-ocr-latest",
        include_image_base64=True
    )

    # Convert response to JSON format
    response_dict = json.loads(pdf_response.model_dump_json())
    all_data = []
    for data in response_dict.get("pages"):
        img=[]

        for image in data.get("images"):
            img.append({
                'id': image.get('id'),
                'desc':image_to_detail(image.get('image_base64'))
            })
        temp={
            'index': data.get('index'),
            'markdown': data.get('markdown'),
            'images': img
        }
        all_data.append(temp)
    print(len(all_data))    
    ppt_data=[]
    for data in range(len(all_data)):
        img = []
        arr = response_dict.get("pages")
        for image in arr[data].get('images'):
            img.append(image.get("image_base64"))
        # all_data[data]['images'] = img
        ppt_data.append([summarize_data_for_ppt(all_data[data]),img])
        print("done")

    create_presentation(ppt_data, "temp\output_presentation.pptx")
    file_path = "test.pdf"
    if os.path.exists(file_path):  # Check if the file exists
        os.remove(file_path)  # Delete the file
        print(f"{file_path} has been deleted.")
    else:
        print("File not found.")
    return "temp\output_presentation.pptx"
    

